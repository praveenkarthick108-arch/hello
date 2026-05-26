from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Palette ───────────────────────────────────────────────────────────────────
DARK_BLUE  = RGBColor(0x1A, 0x3C, 0x5E)
MID_BLUE   = RGBColor(0x2E, 0x6D, 0xA4)
LIGHT_BLUE = RGBColor(0xE8, 0xF4, 0xFF)
ACCENT     = RGBColor(0x00, 0xB0, 0x8A)
ORANGE     = RGBColor(0xF0, 0x8C, 0x00)
RED        = RGBColor(0xCC, 0x00, 0x00)
PURPLE     = RGBColor(0x6A, 0x5A, 0xCD)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GREY  = RGBColor(0x33, 0x33, 0x33)
MID_GREY   = RGBColor(0x66, 0x66, 0x66)
LIGHT_GREY = RGBColor(0xF4, 0xF6, 0xF9)
YELLOW     = RGBColor(0xFF, 0xD7, 0x00)
TEAL       = RGBColor(0x00, 0x89, 0x9E)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.50)
BLANK = prs.slide_layouts[6]

# Layout constants
FOOTER_Y  = 7.08   # footer rect top edge
CONTENT_Y = 1.15   # content starts after header
CONTENT_B = 6.92   # all content must end at or above this line
MARGIN    = 0.25   # left/right outer margin
RIGHT_MAX = 13.08  # rightmost safe x


# ── Primitives ────────────────────────────────────────────────────────────────

def rect(slide, l, t, w, h, fill=None, line=None, lw=Pt(0)):
    sh = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.line.width = lw
    if fill:
        sh.fill.solid(); sh.fill.fore_color.rgb = fill
    else:
        sh.fill.background()
    if line:
        sh.line.color.rgb = line
    else:
        sh.line.fill.background()
    return sh


def txt(slide, text, l, t, w, h, size=13, bold=False, color=DARK_GREY,
        align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tb.word_wrap = True
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tb


def header(slide, title, subtitle=None, bar=DARK_BLUE):
    """Standard slide header — 1.05-inch bar + 0.04-inch accent line."""
    rect(slide, 0, 0, 13.33, 7.50, fill=LIGHT_GREY)
    rect(slide, 0, 0, 13.33, 1.05, fill=bar)
    rect(slide, 0, 1.05, 13.33, 0.04, fill=MID_BLUE)
    txt(slide, title, 0.35, 0.10, 12.6, 0.62, size=26, bold=True, color=WHITE)
    if subtitle:
        txt(slide, subtitle, 0.35, 0.70, 12.6, 0.35, size=12,
            color=RGBColor(0xB8, 0xD4, 0xF0))


def footer(slide, text, bg=DARK_BLUE):
    """Footer bar pinned to y=7.08."""
    rect(slide, 0, FOOTER_Y, 13.33, 0.42, fill=bg)
    txt(slide, "  " + text, 0, FOOTER_Y + 0.02, 13.2, 0.38,
        size=10, color=WHITE, bold=True)


def col_card(slide, l, t, w, h, title, bullets,
             hdr=MID_BLUE, body_size=11, bullet_h=0.36, bullet_pad=0.10):
    """Card with coloured header + bulleted body. Body auto-fits inside h."""
    rect(slide, l, t, w, h, fill=WHITE, line=hdr, lw=Pt(1.5))
    rect(slide, l, t, w, 0.44, fill=hdr)
    txt(slide, title, l + 0.14, t + 0.08, w - 0.26, 0.32,
        size=13, bold=True, color=WHITE)
    body_y = t + 0.44 + bullet_pad
    for b in bullets:
        txt(slide, b, l + 0.16, body_y, w - 0.28, bullet_h,
            size=body_size, color=DARK_GREY)
        body_y += bullet_h + 0.04


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, 13.33, 7.50, fill=DARK_BLUE)
rect(s, 0, 0, 13.33, 0.07, fill=MID_BLUE)
rect(s, 0, 7.43, 13.33, 0.07, fill=MID_BLUE)

# Decorative circles — kept inside RIGHT_MAX = 13.08
for cx, cy, r, c in [
    (10.0, 0.5, 2.2, MID_BLUE),
    (11.5, 2.2, 1.3, RGBColor(0x0D, 0x2A, 0x45)),
    (9.8,  5.8, 1.5, RGBColor(0x0D, 0x2A, 0x45)),
]:
    sh = s.shapes.add_shape(9, Inches(cx), Inches(cy), Inches(r), Inches(r))
    sh.fill.solid(); sh.fill.fore_color.rgb = c; sh.line.fill.background()

txt(s, "AI Trip Planner", 0.7, 1.3, 9.8, 1.3,
    size=54, bold=True, color=WHITE)
txt(s, "Multi-Agent AI System  |  LangGraph  |  GPT-4o-mini",
    0.7, 2.8, 9.5, 0.52, size=18, color=RGBColor(0xB8, 0xD4, 0xF0))
rect(s, 0.7, 3.45, 3.8, 0.05, fill=ACCENT)

for i, b in enumerate([
    "11 specialised AI agents working in parallel",
    "3-layer guardrails — input + output safety",
    "RAGAS + DeepEval quality evaluation suite",
    "Natural language in  ->  Professional PDF out",
]):
    txt(s, f"  {b}", 0.75, 3.65 + i * 0.52, 8.5, 0.46,
        size=15, color=RGBColor(0xCC, 0xE5, 0xFF))

txt(s, "LangGraph  ·  LangChain  ·  OpenAI  ·  Tavily  ·  ChromaDB  ·  RAGAS  ·  DeepEval  ·  FastAPI  ·  Streamlit",
    0.7, 6.65, 12.0, 0.44, size=10, color=RGBColor(0x70, 0x9A, 0xC0))


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — PROBLEM & SOLUTION
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header(s, "Problem & Solution", "What gap does this system fill?")

# Two columns: each 5.9 wide, gap of 1.23 in the middle
for x, bg_, hdr_clr, mark, items in [
    (MARGIN, WHITE, RED, "x", [
        "Hours of research across multiple sites",
        "No single tool: flights + hotels + weather",
        "Manual budget tracking — error-prone",
        "Generic plans — no personalisation",
        "No quality check before output",
        "No safety net for harmful inputs",
    ]),
    (7.18, WHITE, ACCENT, "v", [
        "One natural-language request does it all",
        "11 agents run in parallel for speed",
        "Deterministic budget math — always correct",
        "Semantic memory of past trips",
        "AI quality gate reviews the plan first",
        "Guardrails block harmful / off-topic input",
    ]),
]:
    rect(s, x, 1.18, 5.90, 5.74, fill=bg_, line=hdr_clr, lw=Pt(2))
    rect(s, x, 1.18, 5.90, 0.48, fill=hdr_clr)
    lbl = "  The Problem" if mark == "x" else "  The Solution"
    txt(s, lbl, x, 1.18, 5.90, 0.48, size=15, bold=True, color=WHITE)
    ink = RGBColor(0x88, 0x00, 0x00) if mark == "x" else RGBColor(0x00, 0x70, 0x50)
    for j, p in enumerate(items):
        txt(s, f"  {mark}  {p}", x + 0.18, 1.76 + j * 0.72, 5.55, 0.65,
            size=13, color=ink)

txt(s, "->", 6.17, 3.75, 0.84, 0.72, size=34, bold=True,
    color=MID_BLUE, align=PP_ALIGN.CENTER)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — SYSTEM ARCHITECTURE
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header(s, "System Architecture",
       "11 agents · 3-layer guardrails · evaluation suite")

# ── Row 1: top components (y=1.18) ───────────────────────────────────────────
top_boxes = [
    (0.25, 2.00, "Streamlit UI",        MID_BLUE),
    (2.50, 1.80, "FastAPI",             ORANGE),
    (4.55, 2.55, "Orchestrator\n(LangGraph)", DARK_BLUE),
    (7.35, 1.90, "SQLite DB",           PURPLE),
    (9.50, 2.25, "ChromaDB",            TEAL),
]
for bx, bw, label, bg_ in top_boxes:
    rect(s, bx, 1.18, bw, 0.68, fill=bg_)
    txt(s, label, bx, 1.18, bw, 0.68, size=11, bold=True,
        color=WHITE, align=PP_ALIGN.CENTER)

# arrows between top boxes (small right-arrows)
for ax in [2.30, 4.10, 6.90, 9.05]:
    txt(s, ">", ax, 1.28, 0.30, 0.46, size=16, bold=True,
        color=DARK_BLUE, align=PP_ALIGN.CENTER)

# Guardrail / eval labels below top row
rect(s, 0.25, 1.92, 3.30, 0.26, fill=RED)
txt(s, "  INPUT GUARDRAILS (3 layers)", 0.25, 1.92, 3.30, 0.26,
    size=9, bold=True, color=WHITE)
rect(s, 9.50, 1.92, 3.33, 0.26, fill=PURPLE)
txt(s, "EVAL SUITE (basic+RAGAS+DeepEval)", 9.50, 1.92, 3.33, 0.26,
    size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Down arrow from orchestrator
txt(s, "|", 5.55, 1.90, 0.40, 0.30, size=14, bold=True,
    color=MID_BLUE, align=PP_ALIGN.CENTER)

# ── Row 2: sequential agents phase 0-2 (y=2.28) ──────────────────────────────
txt(s, "Sequential:", 0.25, 2.28, 1.40, 0.28, size=10, bold=True, color=DARK_BLUE)
for i, (label, bg_) in enumerate([
    ("Guardrail\nAgent", RED),
    ("User Input\nAgent", MID_BLUE),
    ("Memory\nAgent", MID_BLUE),
]):
    bx = 0.25 + i * 2.15
    rect(s, bx, 2.58, 2.00, 0.72, fill=bg_)
    txt(s, label, bx, 2.58, 2.00, 0.72, size=11, bold=True,
        color=WHITE, align=PP_ALIGN.CENTER)
    if i < 2:
        txt(s, ">", bx + 2.00, 2.72, 0.15, 0.45, size=14, bold=True,
            color=MID_BLUE, align=PP_ALIGN.CENTER)

# ── Row 3: parallel agents (y=3.42) ──────────────────────────────────────────
txt(s, "Parallel:", 6.65, 2.28, 1.10, 0.28, size=10, bold=True, color=DARK_BLUE)
rect(s, 6.55, 2.50, 6.28, 0.20, fill=RGBColor(0xE0, 0xF0, 0xFF),
     line=MID_BLUE, lw=Pt(1))
txt(s, "  All 4 agents fire simultaneously via Send()", 6.60, 2.50, 6.18, 0.20,
    size=9, bold=True, color=DARK_BLUE)
for i, (label, bg_) in enumerate([
    ("Weather", RGBColor(0x00, 0x7A, 0xCC)),
    ("Transport", RGBColor(0x00, 0x7A, 0xCC)),
    ("Hotel", RGBColor(0x00, 0x7A, 0xCC)),
    ("Places", RGBColor(0x00, 0x7A, 0xCC)),
]):
    bx = 6.55 + i * 1.58
    rect(s, bx, 2.72, 1.46, 0.60, fill=bg_)
    txt(s, label, bx, 2.72, 1.46, 0.60, size=11, bold=True,
        color=WHITE, align=PP_ALIGN.CENTER)

# ── Row 4: bottom sequential (y=3.48) ────────────────────────────────────────
txt(s, "Then sequential:", 0.25, 3.44, 2.10, 0.28, size=10, bold=True, color=DARK_BLUE)
bottom = [
    ("Budget",    ACCENT),
    ("Itinerary", ACCENT),
    ("Review",    ORANGE),
    ("PDF Gen.",  PURPLE),
    ("Output\nGuardrails", RED),
    ("Eval\nSuite", TEAL),
]
# 6 boxes across 12.83 inches: each box = 2.00, gap = 0.14
BW, BG = 2.00, 0.14
for i, (label, bg_) in enumerate(bottom):
    bx = MARGIN + i * (BW + BG)
    rect(s, bx, 3.74, BW, 0.72, fill=bg_)
    txt(s, label, bx, 3.74, BW, 0.72, size=10, bold=True,
        color=WHITE, align=PP_ALIGN.CENTER)
    if i < 5:
        txt(s, ">", bx + BW, 3.90, BG, 0.38, size=11, bold=True,
            color=MID_BLUE, align=PP_ALIGN.CENTER)

# ── Tech strip ────────────────────────────────────────────────────────────────
rect(s, MARGIN, 4.60, 12.83, 0.28, fill=DARK_BLUE)
txt(s, "  LangGraph  |  GPT-4o-mini  |  Tavily  |  ChromaDB  |  SQLite  |  RAGAS  |  DeepEval",
    0.35, 4.61, 12.60, 0.26, size=10, color=WHITE, bold=True)

# ── Key metric tiles ──────────────────────────────────────────────────────────
tiles = [
    ("11", "AI Agents",       MID_BLUE),
    ("3",  "Guard. Layers",   RED),
    ("3",  "Eval Suites",     TEAL),
    ("76", "Tests",           ACCENT),
    ("~90s", "Plan Time",     ORANGE),
]
TW = 2.45
for i, (val, lbl, bg_) in enumerate(tiles):
    bx = MARGIN + i * (TW + 0.11)
    rect(s, bx, 5.00, TW, 1.90, fill=bg_)
    txt(s, val, bx, 5.05, TW, 1.00, size=30, bold=True,
        color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, lbl, bx, 5.95, TW, 0.88, size=11, bold=True,
        color=WHITE, align=PP_ALIGN.CENTER)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — THE 8-PHASE PIPELINE
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header(s, "The 8-Phase Pipeline",
       "Every trip flows through these phases — orchestrator decides what runs next")

phases = [
    ("0", "Guardrail Agent",      "Block harmful/off-topic input before any LLM call", RED),
    ("1", "User Input Agent",     "Parse NL request -> TripPreferences (Pydantic)",     MID_BLUE),
    ("2", "Memory Agent",         "ChromaDB semantic search -> personalise from history",TEAL),
    ("3", "Research (Parallel)",  "Weather + Transport + Hotel + Places — all at once", RGBColor(0x00, 0x7A, 0xCC)),
    ("4", "Budget Agent",         "Deterministic Python math — zero LLM calls",         ACCENT),
    ("5", "Itinerary Agent",      "GPT-4o-mini builds day-by-day schedule",              ACCENT),
    ("6", "Final Review Agent",   "LLM QA gate — detects conflicts, triggers retries",   ORANGE),
    ("7", "PDF Generator",        "ReportLab 8-section report + store to ChromaDB",      PURPLE),
]

# 2 rows of 4 — each card: w=3.10, h=1.72
CW, CH = 3.10, 1.72
COL_GAP = (12.83 - 4 * CW) / 3   # = 0.11
for i, (num, name, desc, bg_) in enumerate(phases):
    col = i % 4
    row = i // 4
    x = MARGIN + col * (CW + COL_GAP)
    y = CONTENT_Y + row * 1.86

    rect(s, x, y, CW, CH, fill=WHITE, line=bg_, lw=Pt(2))
    rect(s, x, y, 0.54, CH, fill=bg_)
    txt(s, num, x, y, 0.54, CH, size=22, bold=True,
        color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, name, x + 0.62, y + 0.12, CW - 0.72, 0.40,
        size=12, bold=True, color=DARK_BLUE)
    txt(s, desc, x + 0.62, y + 0.56, CW - 0.72, 1.06,
        size=11, color=DARK_GREY)

    # connecting arrow (skip last in each row)
    if col < 3:
        txt(s, ">", x + CW + 0.01, y + 0.60, COL_GAP - 0.01, 0.50,
            size=14, bold=True, color=MID_BLUE, align=PP_ALIGN.CENTER)

footer(s, "Orchestrator is pure Python if/else — zero LLM calls — runs after every agent node")


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — INPUT GUARDRAILS
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header(s, "Input Guardrails  —  3 Layers",
       "src/guardrails.py  |  Runs before any LLM call", bar=RED)

# 3 equal columns — w=4.11, h ends at CONTENT_B
CW = 4.11
layers = [
    ("Layer 1", "Deterministic Rules", MID_BLUE, [
        "• Input length: 10–2,000 characters",
        "• Must contain 3+ readable alpha chars",
        "• Dates: not in past, end > start, max 60 days",
        "• Budget $50–$500k   |   Travelers 1–20",
        "• Numeric-only destination blocked",
        "• Zero LLM calls — always fast & free",
    ]),
    ("Layer 2", "LLM Safety Check", ORANGE, [
        "• Is content safe? (no violence, illegal acts)",
        "• Is it travel-related? (lenient — vague OK)",
        "• Prompt injection? (ignore-instructions, DAN…)",
        "• Severity:  pass  /  warn  /  block",
        "• Fail-open: LLM error = pass with warn",
        "• Input truncated to 600 chars for speed",
    ]),
    ("Layer 3", "PII Detection & Masking", TEAL, [
        "• Detects: email, phone, SSN, credit card",
        "• Also: passport number, Aadhaar",
        "• Masks with tokens: [PII_EMAIL_1]",
        "• Restores original values after LLM parsing",
        "• sanitize_for_storage() before any DB write",
        "• Regex-based — no LLM needed",
    ]),
]

CARD_H = CONTENT_B - CONTENT_Y   # 6.92 - 1.15 = 5.77
for i, (num, name, bg_, bullets) in enumerate(layers):
    x = MARGIN + i * (CW + 0.125)
    rect(s, x, CONTENT_Y, CW, CARD_H, fill=WHITE, line=bg_, lw=Pt(2))
    rect(s, x, CONTENT_Y, CW, 0.80, fill=bg_)
    txt(s, num, x + 0.14, CONTENT_Y + 0.04, 0.90, 0.40,
        size=20, bold=True, color=WHITE)
    txt(s, name, x + 0.14, CONTENT_Y + 0.44, CW - 0.22, 0.32,
        size=12, bold=True, color=WHITE)
    for j, b in enumerate(bullets):
        txt(s, b, x + 0.18, CONTENT_Y + 0.92 + j * 0.80, CW - 0.30, 0.74,
            size=11, color=DARK_GREY)

footer(s, "If ANY layer blocks -> request rejected immediately. Blocked = zero LLM cost incurred.")


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — OUTPUT GUARDRAILS
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header(s, "Output Guardrails",
       "src/output_guardrails.py  |  Validates what agents return before downstream use",
       bar=RED)

txt(s, "6 post-generation validators — run after each agent writes to state",
    MARGIN, 1.14, 12.83, 0.32, size=12, bold=True, color=DARK_BLUE)

validators = [
    ("Transport",  RGBColor(0x00, 0x7A, 0xCC),
     "Has outbound options?  |  Cost in range?  |  Not over total budget?"),
    ("Hotel",      ORANGE,
     "Has options?  |  Cost < 80% of budget?  |  Stars rating 1–6?"),
    ("Places",     TEAL,
     "Enough attractions for trip length?  |  2+ restaurants?  |  Indoor options present?"),
    ("Weather",    MID_BLUE,
     "Summary present?  |  Severe weather -> flags for indoor-focused plan"),
    ("Itinerary",  ACCENT,
     "No duplicate dates  |  All dates within trip range  |  No suspiciously short activities"),
    ("Budget",     PURPLE,
     "Components sum = total_spent (within 5%)  |  No negative costs  |  Not 10%+ over budget"),
]

# 2 columns x 3 rows — each card w=6.28, h=1.52
CW, CH = 6.28, 1.52
for i, (name, bg_, check) in enumerate(validators):
    col = i % 2
    row = i // 2
    x = MARGIN + col * (CW + 0.27)
    y = 1.52 + row * (CH + 0.12)
    rect(s, x, y, CW, CH, fill=WHITE, line=bg_, lw=Pt(1.5))
    rect(s, x, y, CW, 0.40, fill=bg_)
    txt(s, f"validate_{name.lower()}_output()", x + 0.14, y + 0.05,
        CW - 0.24, 0.32, size=12, bold=True, color=WHITE)
    txt(s, check, x + 0.16, y + 0.50, CW - 0.28, 0.90, size=11, color=DARK_GREY)

# strip inside content area
rect(s, MARGIN, 6.16, 12.83, 0.36, fill=DARK_BLUE)
txt(s, "  agent_output_guardrails.run_full_output_check(state)  —  runs all 6 checks in one call  ->  {check_name: result}",
    MARGIN + 0.06, 6.18, 12.65, 0.30, size=10, color=WHITE, bold=True)

footer(s, "Severity: pass = ok  |  warn = flagged but continues  |  block = retries agent.  Mirrors input guardrail design.")


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — TECHNOLOGY STACK
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header(s, "Technology Stack", "Every layer of the system")

cats = [
    ("Orchestration",     MID_BLUE, [
        "LangGraph — StateGraph, conditional edges, Send()",
        "LangChain — LLM abstractions & prompt templates",
        "Async / Await throughout for concurrency",
    ]),
    ("AI & LLM",          DARK_BLUE, [
        "OpenAI GPT-4o-mini — all LLM tasks (temp 0–0.3)",
        ".with_structured_output() — typed JSON from LLM",
        "Tavily API — real-time web search",
    ]),
    ("Memory & Storage",  PURPLE, [
        "ChromaDB — vector DB for past trip memory",
        "SQLite + SQLAlchemy — async session tracking",
        "text-embedding-3-small — trip embeddings",
    ]),
    ("Guardrails",        RED, [
        "3-layer input: rules + LLM safety + PII masking",
        "6-check output: per-agent post-gen validation",
        "Regex PII detection — zero extra LLM calls",
    ]),
    ("Evaluation",        TEAL, [
        "RAGAS — Faithfulness + Answer Relevancy",
        "DeepEval — Relevancy, Faithfulness, Hallucination, GEval",
        "Basic Evals — 7 deterministic metrics, 76 tests",
    ]),
    ("Output & API",      ACCENT, [
        "ReportLab — 8-section professional PDF",
        "FastAPI — async REST API (5 endpoints)",
        "Streamlit — live-polling web UI",
    ]),
]

# 3 columns x 2 rows — w=4.11, h=2.60
CW, CH = 4.11, 2.60
for i, (name, bg_, items) in enumerate(cats):
    col = i % 3
    row = i // 3
    x = MARGIN + col * (CW + 0.125)
    y = CONTENT_Y + row * (CH + 0.14)
    rect(s, x, y, CW, CH, fill=WHITE, line=bg_, lw=Pt(2))
    rect(s, x, y, CW, 0.44, fill=bg_)
    txt(s, name, x + 0.14, y + 0.07, CW - 0.22, 0.32,
        size=13, bold=True, color=WHITE)
    for j, item in enumerate(items):
        txt(s, f"• {item}", x + 0.16, y + 0.54 + j * 0.66,
            CW - 0.28, 0.60, size=11, color=DARK_GREY)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — LANGGRAPH STATE & ORCHESTRATOR
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header(s, "LangGraph: Shared State & Orchestrator",
       "src/state.py + src/graph.py  |  Every agent reads from and writes to one TypedDict")

CARD_H = CONTENT_B - CONTENT_Y   # 5.77

# Left: State column
LW = 5.72
rect(s, MARGIN, CONTENT_Y, LW, CARD_H, fill=WHITE, line=MID_BLUE, lw=Pt(2))
rect(s, MARGIN, CONTENT_Y, LW, 0.44, fill=DARK_BLUE)
txt(s, "  TripPlannerState  (TypedDict)", MARGIN, CONTENT_Y, LW, 0.44,
    size=13, bold=True, color=WHITE)

fields = [
    ("user_id, session_id, raw_input", "Set once — immutable",               MID_BLUE),
    ("guardrail_result",               "merge_dicts — phase 0 output",        RED),
    ("trip_preferences",               "merge_dicts — parsed preferences",    ACCENT),
    ("user_profile",                   "merge_dicts — memory lookup result",  TEAL),
    ("weather / transport / hotel / places", "merge_dicts — parallel outputs",RGBColor(0x00, 0x7A, 0xCC)),
    ("budget_summary",                 "merge_dicts — cost breakdown",        ACCENT),
    ("itinerary, review_status",       "merge_dicts — plan + QA result",      ORANGE),
    ("retry_counts",                   "increment_retries — accumulates",     ORANGE),
    ("errors",                         "operator.add — append-only log",      RED),
]
for i, (field, ann, clr) in enumerate(fields):
    y = CONTENT_Y + 0.54 + i * 0.56
    rect(s, MARGIN + 0.05, y, LW - 0.10, 0.50, fill=LIGHT_GREY if i % 2 == 0 else WHITE)
    txt(s, field, MARGIN + 0.14, y + 0.08, 3.00, 0.36,
        size=10, bold=True, color=DARK_BLUE)
    txt(s, ann, MARGIN + 3.20, y + 0.08, LW - 3.30, 0.36,
        size=9, color=clr, italic=True)

# Right: Orchestrator column
RX = MARGIN + LW + 0.28
RW = RIGHT_MAX - RX  # 13.08 - 6.25 = 6.83
rect(s, RX, CONTENT_Y, RW, CARD_H, fill=WHITE, line=DARK_BLUE, lw=Pt(2))
rect(s, RX, CONTENT_Y, RW, 0.44, fill=DARK_BLUE)
txt(s, "  Orchestrator — decision logic (pure Python)", RX, CONTENT_Y, RW, 0.44,
    size=13, bold=True, color=WHITE)

orch_rows = [
    ("0", "no guardrail_result?        -> guardrail_agent",    RED),
    ("1", "no trip_preferences?        -> user_input_agent",   MID_BLUE),
    ("2", "no user_profile?            -> memory_agent",       TEAL),
    ("3", "research missing?           -> Send x4 parallel",   RGBColor(0x00, 0x7A, 0xCC)),
    ("R", "hotel > 60% budget?         -> retry hotel_agent",  ORANGE),
    ("R", "severe weather, no indoor?  -> retry places_agent", ORANGE),
    ("4", "no budget_summary?          -> budget_agent",       ACCENT),
    ("5", "no itinerary?               -> itinerary_agent",    ACCENT),
    ("6", "not reviewed?               -> final_review_agent", ORANGE),
    ("7", "not generated?              -> pdf_generator_agent",PURPLE),
]
for i, (phase, cond, clr) in enumerate(orch_rows):
    y = CONTENT_Y + 0.54 + i * 0.52
    rect(s, RX + 0.05, y, 0.38, 0.44, fill=clr)
    txt(s, phase, RX + 0.05, y, 0.38, 0.44, size=10, bold=True,
        color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, cond, RX + 0.50, y + 0.08, RW - 0.60, 0.34,
        size=10, color=DARK_GREY, italic=True)

footer(s, "Key: returning a LIST of Send() objects = parallel execution.  Returning a string = single next node.")


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — PARALLEL RESEARCH + RETRY LOGIC
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header(s, "Parallel Research  +  Retry Logic",
       "4 agents simultaneously · Orchestrator detects problems and retries the right one")

# Section label
rect(s, MARGIN, 1.14, 12.83, 0.30, fill=RGBColor(0x00, 0x7A, 0xCC))
txt(s, "  Phase 3  —  All 4 fired via Send() at the same moment  —  saves ~45 seconds vs sequential",
    MARGIN + 0.06, 1.15, 12.65, 0.28, size=11, bold=True, color=WHITE)

# 4 agent cards — w=3.10, h=2.52
AW, AH = 3.10, 2.52
for i, (name, tool, outs) in enumerate([
    ("Weather Agent",   "OpenWeatherMap API (mock if no key)",
     ["5-day daily forecast", "severe_weather flag", "rain_days count", "summary string"]),
    ("Transport Agent", "Tavily search + LLM parse",
     ["Outbound + return options", "Recommended option", "estimated_cost_usd", "availability flag"]),
    ("Hotel Agent",     "Tavily search + LLM parse",
     ["Hotel options list", "Recommended hotel dict", "estimated_cost_usd", "within_budget flag"]),
    ("Places Agent",    "Tavily x3 searches",
     ["Attractions list", "Restaurants list", "Local experiences", "Indoor options (rain backup)"]),
]):
    x = MARGIN + i * (AW + 0.05)
    y = 1.50
    rect(s, x, y, AW, AH, fill=WHITE, line=RGBColor(0x00, 0x7A, 0xCC), lw=Pt(1.5))
    rect(s, x, y, AW, 0.40, fill=RGBColor(0x00, 0x7A, 0xCC))
    txt(s, name, x + 0.12, y + 0.06, AW - 0.20, 0.30, size=12, bold=True, color=WHITE)
    rect(s, x, y + 0.40, AW, 0.30, fill=LIGHT_GREY)
    txt(s, tool, x + 0.12, y + 0.42, AW - 0.20, 0.26,
        size=9, color=MID_GREY, italic=True)
    for j, o in enumerate(outs):
        txt(s, f"-> {o}", x + 0.12, y + 0.80 + j * 0.44, AW - 0.20, 0.38,
            size=10, color=DARK_GREY)

# Retry section
rect(s, MARGIN, 4.15, 12.83, 0.30, fill=ORANGE)
txt(s, "  Retry Logic  —  Auto-retries specific agents when conflicts detected (max 3 each)",
    MARGIN + 0.06, 4.16, 12.65, 0.28, size=11, bold=True, color=WHITE)

retries = [
    ("Hotel over budget",     "hotel_cost > budget x 0.60  and  not within_budget",
     "Retry hotel_agent  (searches for cheaper options)"),
    ("Transport unavailable", "transport.availability = False",
     "Retry transport_agent  (tries alternate routes)"),
    ("Severe weather",        "severe_weather = True  and  no indoor_options",
     "Retry places_agent  (searches indoor venues)"),
    ("Review rejected",       "review.approved = False  and  not _retried",
     "Retry listed agent  ->  re-run review agent"),
]

RW2 = (12.83 - 3 * 0.10) / 4   # = 3.13
for i, (title, cond, result) in enumerate(retries):
    x = MARGIN + i * (RW2 + 0.10)
    y = 4.52
    rect(s, x, y, RW2, 2.30, fill=WHITE, line=ORANGE, lw=Pt(1.5))
    rect(s, x, y, RW2, 0.38, fill=ORANGE)
    txt(s, title, x + 0.12, y + 0.06, RW2 - 0.20, 0.28, size=11, bold=True, color=WHITE)
    rect(s, x, y + 0.38, RW2, 0.88, fill=RGBColor(0xFF, 0xF5, 0xE6))
    txt(s, f"IF:  {cond}", x + 0.12, y + 0.42, RW2 - 0.20, 0.78,
        size=9.5, color=DARK_GREY, italic=True)
    rect(s, x, y + 1.26, RW2, 0.90, fill=RGBColor(0xF0, 0xFF, 0xF0))
    txt(s, f"THEN:  {result}", x + 0.12, y + 1.30, RW2 - 0.20, 0.80,
        size=10, color=RGBColor(0x00, 0x60, 0x00))

footer(s, "retry_counts uses increment_retries reducer — adds not overwrites, so limits are correctly enforced")


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — BUDGET AGENT + FINAL REVIEW AGENT
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header(s, "Budget Agent  +  Final Review Agent",
       "Deterministic math then AI quality gate before the PDF is generated")

CARD_H = CONTENT_B - CONTENT_Y   # 5.77
LW = 5.88
RX = MARGIN + LW + 0.30
RW = RIGHT_MAX - RX   # 6.65

# ── Budget left ──────────────────────────────────────────────────────────────
rect(s, MARGIN, CONTENT_Y, LW, CARD_H, fill=WHITE, line=ACCENT, lw=Pt(2))
rect(s, MARGIN, CONTENT_Y, LW, 0.44, fill=ACCENT)
txt(s, "  Budget Agent  —  Zero LLM Calls", MARGIN, CONTENT_Y, LW, 0.44,
    size=14, bold=True, color=WHITE)
txt(s, "Why? Costs are arithmetic, not language. LLMs hallucinate numbers — Python math is exact.",
    MARGIN + 0.16, CONTENT_Y + 0.52, LW - 0.26, 0.44, size=11, color=MID_GREY, italic=True)

calcs = [
    ("transport_cost",  "transport.estimated_cost_usd  (fallback: budget x 0.25)"),
    ("hotel_cost",      "hotel.estimated_cost_usd  (fallback: budget x 0.40)"),
    ("activities_cost", "top-5 attraction fees x travelers"),
    ("food_cost",       "days x travelers x $40 / day flat rate"),
    ("misc_cost",       "total_budget x 0.05 contingency"),
    ("total_spent",     "sum of all components above"),
    ("remaining",       "total_budget - total_spent"),
    ("within_budget",   "total_spent <= total_budget"),
]
for i, (field, formula) in enumerate(calcs):
    y = CONTENT_Y + 1.04 + i * 0.55
    rect(s, MARGIN + 0.05, y, LW - 0.10, 0.50,
         fill=LIGHT_GREY if i % 2 == 0 else WHITE)
    txt(s, field, MARGIN + 0.14, y + 0.08, 1.90, 0.36,
        size=11, bold=True, color=DARK_BLUE)
    txt(s, formula, MARGIN + 2.10, y + 0.08, LW - 2.22, 0.36,
        size=10, color=DARK_GREY, italic=True)

# ── Final Review right ───────────────────────────────────────────────────────
rect(s, RX, CONTENT_Y, RW, CARD_H, fill=WHITE, line=ORANGE, lw=Pt(2))
rect(s, RX, CONTENT_Y, RW, 0.44, fill=ORANGE)
txt(s, "  Final Review Agent  —  AI Quality Gate", RX, CONTENT_Y, RW, 0.44,
    size=14, bold=True, color=WHITE)

checks = [
    ("Budget compliance",    "total_spent > budget by >10%?"),
    ("Date consistency",     "itinerary days match start->end dates?"),
    ("Hotel-transport fit",  "hotel check-in = arrival day?"),
    ("Weather fit",          "outdoor activities on heavy rain days?"),
    ("Completeness",         "every day has at least one activity?"),
]
for i, (name, q) in enumerate(checks):
    y = CONTENT_Y + 0.54 + i * 0.62
    rect(s, RX + 0.05, y, 0.40, 0.54, fill=RED)
    txt(s, str(i + 1), RX + 0.05, y, 0.40, 0.54, size=13, bold=True,
        color=WHITE, align=PP_ALIGN.CENTER)
    rect(s, RX + 0.48, y, RW - 0.54, 0.54, fill=WHITE, line=RED, lw=Pt(0.75))
    txt(s, name, RX + 0.58, y + 0.06, RW - 0.78, 0.22, size=11, bold=True, color=DARK_BLUE)
    txt(s, q,   RX + 0.58, y + 0.30, RW - 0.78, 0.22, size=10, color=DARK_GREY, italic=True)

# ReviewResult box — sits below the 5 checks
RR_Y = CONTENT_Y + 0.54 + 5 * 0.62   # = 1.15 + 0.54 + 3.10 = 4.79
RR_H = CONTENT_B - RR_Y - 0.05        # remaining space
rect(s, RX + 0.05, RR_Y, RW - 0.10, RR_H, fill=LIGHT_GREY, line=ORANGE, lw=Pt(1))
txt(s, "ReviewResult output:", RX + 0.18, RR_Y + 0.08, RW - 0.30, 0.28,
    size=11, bold=True, color=DARK_BLUE)
for i, line in enumerate([
    "approved: True only if no blocking conflicts",
    "conflicts: list of blocking issues found",
    "retry_reasons: which agent to re-run",
    "approved=False -> retry once -> PDF anyway",
]):
    txt(s, f"• {line}", RX + 0.18, RR_Y + 0.40 + i * 0.38, RW - 0.30, 0.34,
        size=10, color=DARK_GREY)

footer(s, "Graceful degradation: if review agent crashes -> approved=True with warning — user always gets a plan")


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — EVALUATION FRAMEWORK OVERVIEW
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header(s, "Evaluation Framework",
       "evals/  —  Three suites that measure quality at different levels", bar=TEAL)

txt(s, "Every trip plan is scored across three independent evaluation suites",
    MARGIN, 1.14, 12.83, 0.30, size=12, bold=True, color=DARK_BLUE)

# 3 equal columns — card ends at CONTENT_B
CW = 4.11
suites = [
    ("Basic Evals", ACCENT, "evals/basic_evals.py",
     "7 deterministic checks — no LLM, no API cost",
     [
         "Itinerary completeness (days + slots filled)",
         "Budget arithmetic accuracy",
         "Date consistency across all state fields",
         "Content quality (activity density, coverage)",
         "Guardrail compliance score",
         "Review approval status",
         "Output guardrail summary score",
     ],
     "Speed: instant  |  Cost: free  |  Always reliable"),
    ("RAGAS", TEAL, "evals/ragas_evals.py",
     "RAG pipeline quality — itinerary grounded in research?",
     [
         "Faithfulness: plan grounded in retrieved context?",
         "Answer Relevancy: plan matches user request?",
         "Context = weather + hotel + transport + places",
         "Question = raw_input from user",
         "Answer = full itinerary text",
         "Batch evaluation across multiple trips supported",
     ],
     "Speed: ~30s  |  Cost: OpenAI tokens  |  pip install ragas"),
    ("DeepEval", PURPLE, "evals/deepeval_evals.py",
     "LLM output quality — 4 metrics + custom rubric",
     [
         "AnswerRelevancyMetric — relevant to request?",
         "FaithfulnessMetric — grounded in context?",
         "HallucinationMetric — invented facts?",
         "GEval (custom) — 7-step trip quality rubric",
         "Bonus: ToxicityMetric + BiasMetric on input",
         "Score + reason returned per metric",
     ],
     "Speed: ~60s  |  Cost: OpenAI tokens  |  pip install deepeval"),
]

CARD_H = CONTENT_B - CONTENT_Y  # 5.77
for i, (name, bg_, path, tagline, bullets, foot_txt) in enumerate(suites):
    x = MARGIN + i * (CW + 0.125)
    y = CONTENT_Y
    rect(s, x, y, CW, CARD_H, fill=WHITE, line=bg_, lw=Pt(2))
    rect(s, x, y, CW, 0.76, fill=bg_)
    txt(s, name, x + 0.14, y + 0.05, CW - 0.22, 0.40, size=18, bold=True, color=WHITE)
    txt(s, path, x + 0.14, y + 0.46, CW - 0.22, 0.26, size=9, color=WHITE, italic=True)
    rect(s, x, y + 0.76, CW, 0.28, fill=LIGHT_GREY)
    txt(s, tagline, x + 0.14, y + 0.78, CW - 0.22, 0.24, size=9, bold=True, color=DARK_BLUE)
    for j, b in enumerate(bullets):
        txt(s, f"• {b}", x + 0.16, y + 1.12 + j * 0.46, CW - 0.28, 0.42,
            size=10, color=DARK_GREY)
    # footer strip at bottom of card, inside CONTENT_B
    STRIP_H = 0.30
    rect(s, x, y + CARD_H - STRIP_H, CW, STRIP_H, fill=bg_)
    txt(s, foot_txt, x + 0.10, y + CARD_H - STRIP_H + 0.04,
        CW - 0.16, STRIP_H - 0.08, size=8.5, color=WHITE, bold=True)

footer(s, "EvaluationRunner orchestrates all 3 suites, saves JSON report, prints human-readable summary")


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — BASIC EVALS DETAIL
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header(s, "Basic Evals  —  7 Deterministic Metrics",
       "evals/basic_evals.py  |  Instant, free, always reliable", bar=ACCENT)

metrics = [
    ("Itinerary Completeness", "0.80",
     "days covered / expected (50%)  +  avg time slots filled per day (50%)"),
    ("Budget Accuracy",        "0.90",
     "Component sum == total_spent within 5%  |  No 10%+ budget overrun"),
    ("Date Consistency",       "1.00",
     "start < end, not in past, all itinerary day-dates within trip range"),
    ("Content Quality",        "0.70",
     "Activity density per day, hotel coverage, meal coverage, research data richness"),
    ("Guardrail Compliance",   "1.00",
     "Maps severity to score:  pass=1.0  |  warn=0.7  |  block=0.0"),
    ("Review Approval",        "1.00",
     "approved + no conflicts=1.0  |  approved + conflicts=0.7  |  rejected=0.0"),
    ("Output Guardrails",      "0.80",
     "Averages all 6 output guardrail severity scores into a single 0-1 metric"),
]

# Layout: 2 cols for first 6 (3 rows), metric 7 centered in a full-width row
CW_H = 6.20   # half-width card
CW_F = 12.83  # full-width card
CH = 1.18

for i in range(6):
    col = i % 2
    row = i // 2
    x = MARGIN + col * (CW_H + 0.43)
    y = CONTENT_Y + row * (CH + 0.08)
    name, threshold, desc = metrics[i]
    rect(s, x, y, CW_H, CH, fill=WHITE, line=ACCENT, lw=Pt(1.5))
    rect(s, x, y, 0.52, CH, fill=ACCENT)
    txt(s, str(i + 1), x, y, 0.52, CH, size=20, bold=True,
        color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, name, x + 0.60, y + 0.08, 3.80, 0.36, size=12, bold=True, color=DARK_BLUE)
    txt(s, f"Threshold: {threshold}", x + 4.44, y + 0.10, 1.62, 0.28,
        size=10, bold=True, color=TEAL)
    txt(s, desc, x + 0.60, y + 0.50, CW_H - 0.70, 0.58, size=10, color=DARK_GREY)

# Metric 7 — full width, centered
y7 = CONTENT_Y + 3 * (CH + 0.08)   # row 3
name, threshold, desc = metrics[6]
rect(s, MARGIN, y7, CW_F, CH, fill=WHITE, line=ACCENT, lw=Pt(1.5))
rect(s, MARGIN, y7, 0.52, CH, fill=ACCENT)
txt(s, "7", MARGIN, y7, 0.52, CH, size=20, bold=True,
    color=WHITE, align=PP_ALIGN.CENTER)
txt(s, name, MARGIN + 0.60, y7 + 0.08, 4.50, 0.36, size=12, bold=True, color=DARK_BLUE)
txt(s, f"Threshold: {threshold}", MARGIN + 5.20, y7 + 0.10, 1.62, 0.28,
    size=10, bold=True, color=TEAL)
txt(s, desc, MARGIN + 0.60, y7 + 0.50, CW_F - 0.70, 0.58, size=10, color=DARK_GREY)

# Stat strip — inside content area
STRIP_Y = y7 + CH + 0.10
rect(s, MARGIN, STRIP_Y, 12.83, 0.38, fill=DARK_BLUE)
txt(s, "  BasicEvalSuite.run(state)  ->  {suite, passed, total, pass_rate, avg_score, results}",
    MARGIN + 0.06, STRIP_Y + 0.04, 8.20, 0.30, size=11, color=WHITE, bold=True)
txt(s, "  76 pytest tests — all passing", 8.70, STRIP_Y + 0.04, 4.10, 0.30,
    size=11, color=YELLOW, bold=True)

footer(s, "Runs in <1 second, zero API cost — ideal for CI pipelines and rapid regression checks")


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — RAGAS EVALUATION
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header(s, "RAGAS Evaluation",
       "evals/ragas_evals.py  |  Measures RAG pipeline quality", bar=TEAL)

CARD_H = CONTENT_B - CONTENT_Y   # 5.77
LW = 5.90

# Left: state -> RAGAS mapping
rect(s, MARGIN, CONTENT_Y, LW, CARD_H, fill=WHITE, line=TEAL, lw=Pt(2))
rect(s, MARGIN, CONTENT_Y, LW, 0.44, fill=TEAL)
txt(s, "  State -> RAGAS Field Mapping", MARGIN, CONTENT_Y, LW, 0.44,
    size=13, bold=True, color=WHITE)

mappings = [
    ("user_input",         "raw_input",
     '"Plan a 3-day trip to Tokyo for 2 people\nwith $4000 budget. We enjoy street food."'),
    ("response",           "itinerary text",
     "Day 1 (2026-06-01): Morning: Arrive at Narita...\nBudget: Transport $900, Hotel $1000..."),
    ("retrieved_contexts", "all research outputs",
     "Weather: Warm, 72°F sunny  |  Hotels: Hotel Shinjuku $200/night\nTransport: JAL flight $900  |  Attractions: Senso-ji Temple..."),
]
for i, (ragas_field, source, example) in enumerate(mappings):
    y = CONTENT_Y + 0.54 + i * 1.72
    rect(s, MARGIN + 0.05, y, LW - 0.10, 1.62, fill=LIGHT_GREY if i % 2 == 0 else WHITE,
         line=TEAL, lw=Pt(0.75))
    rect(s, MARGIN + 0.05, y, 1.55, 0.34, fill=TEAL)
    txt(s, ragas_field, MARGIN + 0.05, y, 1.55, 0.34, size=10, bold=True,
        color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, f"Source: state['{source}']", MARGIN + 1.68, y + 0.06, LW - 1.82, 0.26,
        size=10, color=MID_BLUE, bold=True)
    txt(s, example, MARGIN + 0.18, y + 0.42, LW - 0.28, 1.12,
        size=10, color=MID_GREY, italic=True)

# Right: metrics + code
RX = MARGIN + LW + 0.28
RW = RIGHT_MAX - RX   # ~6.65

# Metrics section
rect(s, RX, CONTENT_Y, RW, 2.72, fill=WHITE, line=TEAL, lw=Pt(2))
rect(s, RX, CONTENT_Y, RW, 0.44, fill=TEAL)
txt(s, "  Metrics Used", RX, CONTENT_Y, RW, 0.44, size=13, bold=True, color=WHITE)
for i, (name, what) in enumerate([
    ("Faithfulness",
     "Is every fact in the itinerary grounded in the retrieved contexts?\nE.g. invents a hotel not in hotel_data -> low score."),
    ("Answer Relevancy",
     "Does the itinerary address the user's actual request?\nE.g. asked for Tokyo, got Paris plan -> low score."),
]):
    y = CONTENT_Y + 0.52 + i * 1.06
    rect(s, RX + 0.06, y, RW - 0.12, 0.96, fill=LIGHT_BLUE, line=TEAL, lw=Pt(0.75))
    txt(s, name, RX + 0.18, y + 0.08, RW - 0.30, 0.28, size=12, bold=True, color=DARK_BLUE)
    txt(s, what, RX + 0.18, y + 0.40, RW - 0.30, 0.50, size=10, color=DARK_GREY)

# Code section
CODE_Y = CONTENT_Y + 2.84
CODE_H = CONTENT_B - CODE_Y   # remaining
rect(s, RX, CODE_Y, RW, CODE_H, fill=WHITE, line=MID_BLUE, lw=Pt(2))
rect(s, RX, CODE_Y, RW, 0.40, fill=MID_BLUE)
txt(s, "  Usage — 3 lines", RX, CODE_Y, RW, 0.40, size=12, bold=True, color=WHITE)
txt(s,
    "sample = SingleTurnSample(\n"
    "  user_input = raw_input,\n"
    "  response = itinerary_text,\n"
    "  retrieved_contexts = context_list)\n\n"
    "result = evaluate(\n"
    "  dataset = EvaluationDataset([sample]),\n"
    "  metrics = [Faithfulness(), AnswerRelevancy()])",
    RX + 0.16, CODE_Y + 0.48, RW - 0.24, CODE_H - 0.56,
    size=10, color=DARK_GREY)

footer(s, "evaluate_batch(states) runs RAGAS on multiple trips in one call — ideal for regression testing after prompt changes")


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — DEEPEVAL EVALUATION
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header(s, "DeepEval Evaluation",
       "evals/deepeval_evals.py  |  4 LLM-as-judge metrics + custom rubric", bar=PURPLE)

# 4 metric cards — top row
AW = (12.83 - 3 * 0.10) / 4   # = 3.13
MH = 2.52
for i, (name, bg_, what, low_if) in enumerate([
    ("Answer Relevancy", RGBColor(0x00, 0x7A, 0xCC),
     "Is the itinerary relevant to the user's request?",
     "Plan ignores stated preferences or destination"),
    ("Faithfulness", TEAL,
     "Is every fact supported by retrieved context?",
     "Hotel name not in hotel_data, invented attractions"),
    ("Hallucination", RED,
     "Does itinerary contain invented facts?",
     "Fake hotel names, wrong prices, wrong dates"),
    ("GEval (Custom)", PURPLE,
     "7-step trip quality rubric scored by GPT-4o-mini",
     "Destination mismatch, missing days, vague activities"),
]):
    x = MARGIN + i * (AW + 0.10)
    y = CONTENT_Y
    rect(s, x, y, AW, MH, fill=WHITE, line=bg_, lw=Pt(2))
    rect(s, x, y, AW, 0.40, fill=bg_)
    txt(s, name, x + 0.12, y + 0.06, AW - 0.20, 0.30, size=12, bold=True, color=WHITE)
    txt(s, what, x + 0.12, y + 0.48, AW - 0.20, 0.84, size=11, color=DARK_GREY)
    rect(s, x, y + 1.38, AW, 0.22, fill=LIGHT_GREY)
    txt(s, "Score drops if:", x + 0.12, y + 1.40, 1.50, 0.18, size=8, bold=True, color=MID_GREY)
    txt(s, low_if, x + 0.12, y + 1.62, AW - 0.20, 0.78, size=10, color=MID_GREY, italic=True)

# Bottom left: GEval steps
BH = CONTENT_B - (CONTENT_Y + MH + 0.12)   # = 6.92 - (1.15+2.52+0.12) = 3.13
BLW = 6.20
BY = CONTENT_Y + MH + 0.12
rect(s, MARGIN, BY, BLW, BH, fill=WHITE, line=PURPLE, lw=Pt(2))
rect(s, MARGIN, BY, BLW, 0.40, fill=PURPLE)
txt(s, "  GEval — Custom 7-Step Evaluation Rubric", MARGIN, BY, BLW, 0.40,
    size=13, bold=True, color=WHITE)
steps = [
    "1. Destination matches user's requested destination",
    "2. Itinerary covers the full trip duration (all days)",
    "3. Daily activities are specific — named venues, not vague",
    "4. Accommodation mentioned for each night",
    "5. Budget estimate does not exceed stated limit",
    "6. Daily pace is realistic — not overloaded",
    "7. Transport notes on arrival and departure days",
]
for i, step in enumerate(steps):
    txt(s, step, MARGIN + 0.18, BY + 0.48 + i * 0.38, BLW - 0.28, 0.34,
        size=11, color=DARK_GREY)

# Bottom right: usage
BRX = MARGIN + BLW + 0.28
BRW = RIGHT_MAX - BRX
rect(s, BRX, BY, BRW, BH, fill=WHITE, line=MID_BLUE, lw=Pt(2))
rect(s, BRX, BY, BRW, 0.40, fill=MID_BLUE)
txt(s, "  Usage", BRX, BY, BRW, 0.40, size=13, bold=True, color=WHITE)
txt(s,
    "suite = DeepEvalSuite(\n"
    "  model='gpt-4o-mini',\n"
    "  threshold=0.7)\n\n"
    "result = suite.evaluate(state)\n\n"
    "Returns:\n"
    "  passed:    n metrics above threshold\n"
    "  metrics:   {name: {score, passed, reason}}\n"
    "  avg_score: float 0.0–1.0\n\n"
    "Also available:\n"
    "  evaluate_guardrail(state)\n"
    "  -> ToxicityMetric + BiasMetric\n"
    "     on the raw user input",
    BRX + 0.16, BY + 0.48, BRW - 0.24, BH - 0.56, size=10, color=DARK_GREY)

footer(s, "Each metric calls metric.measure(test_case) individually — easier to catch per-metric errors than the batch evaluate() call")


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — EVAL RESULTS (LIVE RUN)
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header(s, "Eval Results  —  Live Tokyo Trip Test",
       "What the eval suite found on a real plan generated during the demo", bar=TEAL)

txt(s,
    'Input: "Plan a 3-day trip to Tokyo for 2 from New York with $4000 budget. We enjoy street food and anime culture."',
    MARGIN, 1.14, 12.83, 0.30, size=11, bold=True, color=DARK_BLUE)

# Table header
rect(s, MARGIN, 1.50, 7.90, 0.36, fill=DARK_BLUE)
for cx, cw, label in [
    (MARGIN, 2.60, "Evaluator"),
    (MARGIN + 2.60, 1.00, "Score"),
    (MARGIN + 3.60, 0.95, "Pass?"),
    (MARGIN + 4.55, 3.35, "Finding"),
]:
    txt(s, label, cx + 0.08, 1.52, cw - 0.12, 0.28,
        size=11, bold=True, color=WHITE)

rows = [
    ("Itinerary Completeness", "1.000", "PASS",  ACCENT,
     "3/3 days, all 3 time slots filled per day"),
    ("Budget Accuracy",        "1.000", "PASS",  ACCENT,
     "Components sum correctly to total_spent"),
    ("Date Consistency",       "1.000", "PASS",  ACCENT,
     "All dates valid and within trip range"),
    ("Content Quality",        "0.889", "PASS",  ACCENT,
     "Rich activities, hotels, meals present"),
    ("Guardrail Compliance",   "1.000", "PASS",  ACCENT,
     "Input passed all 3 layers cleanly"),
    ("Review Approval",        "0.000", "FAIL",  RED,
     "Review flagged 2% overrun as '>10%' — pre-existing bug!"),
    ("Output Guardrails",      "0.450", "FAIL",  RED,
     "Status API omits research fields — not a pipeline bug"),
]
for i, (name, score, status, clr, note) in enumerate(rows):
    y = 1.90 + i * 0.54
    bg_ = LIGHT_GREY if i % 2 == 0 else WHITE
    rect(s, MARGIN, y, 7.90, 0.50, fill=bg_)
    txt(s, name,   MARGIN + 0.10, y + 0.09, 2.46, 0.32, size=10, bold=True, color=DARK_BLUE)
    txt(s, score,  MARGIN + 2.65, y + 0.09, 0.86, 0.32, size=11, bold=True,
        color=clr, align=PP_ALIGN.CENTER)
    rect(s, MARGIN + 3.60, y + 0.06, 0.88, 0.38, fill=clr)
    txt(s, status, MARGIN + 3.60, y + 0.06, 0.88, 0.38, size=10, bold=True,
        color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, note,   MARGIN + 4.60, y + 0.09, 3.20, 0.32, size=10, color=DARK_GREY)

# Summary strip under table
SUM_Y = 1.90 + 7 * 0.54   # = 5.68
rect(s, MARGIN, SUM_Y, 7.90, 0.36, fill=DARK_BLUE)
txt(s, "  Pass rate: 71.4%  |  Avg score: 0.763  |  5 PASS  2 FAIL  |  Ran in 2 seconds  |  Zero API cost",
    MARGIN + 0.06, SUM_Y + 0.04, 7.72, 0.28, size=10, color=WHITE, bold=True)

# Right: insights panel
RX = MARGIN + 7.90 + 0.28
RW = RIGHT_MAX - RX   # ~4.65
rect(s, RX, 1.50, RW, CONTENT_B - 1.50, fill=WHITE, line=ORANGE, lw=Pt(2))
rect(s, RX, 1.50, RW, 0.40, fill=ORANGE)
txt(s, "  Key Insights", RX, 1.50, RW, 0.40, size=13, bold=True, color=WHITE)

for i, (title, clr, body) in enumerate([
    ("Bug caught by evals", RED,
     "review_approval FAIL exposed that the final_review_agent flags a 2% budget overrun as '>10%' — the threshold in the review prompt is wrong."),
    ("Status API gap", MID_BLUE,
     "output_guardrails FAIL because GET /trips/{id}/status omits research fields. Evals work perfectly on the in-pipeline state object."),
    ("5/7 is real signal", ACCENT,
     "Even basic evals (free, instant) surface real issues in 2 seconds. Immediate value before enabling RAGAS or DeepEval."),
]):
    y = 1.98 + i * 1.58
    rect(s, RX + 0.06, y, RW - 0.12, 0.30, fill=clr)
    txt(s, title, RX + 0.14, y + 0.05, RW - 0.24, 0.22, size=10, bold=True, color=WHITE)
    txt(s, body, RX + 0.14, y + 0.38, RW - 0.24, 1.10, size=10, color=DARK_GREY)

footer(s, "RAGAS and DeepEval (requires OpenAI key) add LLM-graded scores on top of the free deterministic baseline")


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 16 — FASTAPI + STORAGE
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header(s, "FastAPI Backend  +  Storage",
       "api/main.py  |  5 async endpoints  |  SQLite + ChromaDB")

# Endpoints table — left 6.80 wide
TW = 6.80
rect(s, MARGIN, CONTENT_Y, TW, 0.36, fill=DARK_BLUE)
for cx, cw, label in [
    (MARGIN,        1.00, "Method"),
    (MARGIN + 1.00, 3.20, "Endpoint"),
    (MARGIN + 4.20, 0.95, "Status"),
    (MARGIN + 5.15, 1.65, "Returns"),
]:
    txt(s, label, cx + 0.08, CONTENT_Y + 0.04, cw - 0.12, 0.28,
        size=11, bold=True, color=WHITE)

endpoints = [
    ("POST", "/trips",               "202", MID_BLUE, "session_id — graph starts in background"),
    ("GET",  "/trips/{id}/status",   "200", ACCENT,   "current_phase + all results"),
    ("GET",  "/trips/{id}/pdf",      "200", PURPLE,   "FileResponse — PDF download"),
    ("GET",  "/users/{id}/trips",    "200", ORANGE,   "Trip history list"),
    ("GET",  "/health",              "200", DARK_GREY,"version + model name"),
]
for i, (method, path, code, clr, ret) in enumerate(endpoints):
    y = CONTENT_Y + 0.36 + i * 0.50
    rect(s, MARGIN, y, TW, 0.46, fill=LIGHT_GREY if i % 2 == 0 else WHITE)
    rect(s, MARGIN, y, 0.95, 0.46, fill=clr)
    txt(s, method, MARGIN, y, 0.95, 0.46, size=10, bold=True,
        color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, path,   MARGIN + 1.00, y + 0.10, 3.10, 0.28, size=10, bold=True, color=DARK_BLUE)
    rect(s, MARGIN + 4.20, y + 0.05, 0.88, 0.36, fill=ACCENT)
    txt(s, code, MARGIN + 4.20, y + 0.05, 0.88, 0.36, size=10, bold=True,
        color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, ret,  MARGIN + 5.18, y + 0.10, 1.60, 0.28, size=10, color=DARK_GREY)

# SQLite card — below endpoints, ends at CONTENT_B
SQL_Y = CONTENT_Y + 0.36 + 5 * 0.50 + 0.12   # = 1.15 + 0.36 + 2.5 + 0.12 = 4.13
SQL_H = CONTENT_B - SQL_Y   # = 6.92 - 4.13 = 2.79
rect(s, MARGIN, SQL_Y, TW, SQL_H, fill=WHITE, line=PURPLE, lw=Pt(1.5))
rect(s, MARGIN, SQL_Y, TW, 0.40, fill=PURPLE)
txt(s, "  SQLite  —  Session Tracking  (sqlite_manager.py)",
    MARGIN, SQL_Y, TW, 0.40, size=12, bold=True, color=WHITE)
for i, (col, desc) in enumerate([
    ("session_id",     "UUID primary key"),
    ("status",         "in_progress / completed / failed"),
    ("current_phase",  "live phase name — updated after each agent"),
    ("trip_preferences / budget_summary / itinerary", "JSON blobs"),
    ("pdf_path  /  errors",    "PDF file path + append-only error log"),
]):
    txt(s, f"• {col}:  {desc}", MARGIN + 0.16,
        SQL_Y + 0.48 + i * 0.44, TW - 0.26, 0.38, size=10, color=DARK_GREY)

# ChromaDB card — right column
RX = MARGIN + TW + 0.28
RW = RIGHT_MAX - RX
rect(s, RX, CONTENT_Y, RW, CONTENT_B - CONTENT_Y, fill=WHITE, line=TEAL, lw=Pt(2))
rect(s, RX, CONTENT_Y, RW, 0.44, fill=TEAL)
txt(s, "  ChromaDB  —  Semantic Trip Memory", RX, CONTENT_Y, RW, 0.44,
    size=13, bold=True, color=WHITE)

for i, (title, body) in enumerate([
    ("What is stored",   "After each trip, PDF Generator saves a summary as a vector embedding (text-embedding-3-small)."),
    ("Document content", "Destination, dates, budget, activities, preferences. Metadata: user_id, session_id, timestamp."),
    ("How retrieved",    "Memory Agent does semantic search: 'trips to Tokyo ~$4000' -> top-K similar past trips."),
    ("What it enables",  "Returning users get personalised recommendations based on actual past trip preferences."),
    ("Graceful failure", "ChromaDB errors don't block planning. Sets new_user=True and continues without memory."),
]):
    y = CONTENT_Y + 0.54 + i * 1.06
    txt(s, title, RX + 0.16, y, RW - 0.26, 0.28, size=11, bold=True, color=DARK_BLUE)
    txt(s, body,  RX + 0.16, y + 0.30, RW - 0.26, 0.68, size=10, color=DARK_GREY)
    if i < 4:
        rect(s, RX, y + 1.00, RW, 0.04, fill=LIGHT_GREY)

footer(s, "Live progress: graph.astream() yields node-by-node -> update_trip_record(current_phase=node) -> Streamlit sees it in real time")


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 17 — END-TO-END FLOW
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header(s, "End-to-End Flow",
       "From user input to evaluated PDF — 60 to 120 seconds total")

steps = [
    ("1",  "User fills form",         "Origin, destination, dates, budget, preferences",              MID_BLUE),
    ("2",  "POST /trips",             "FastAPI saves to SQLite, fires background task, returns session_id", ORANGE),
    ("3",  "Input Guardrails",        "Deterministic rules  +  LLM safety  +  PII masking",           RED),
    ("4",  "Parsing + Memory",        "user_input_agent -> preferences  |  memory_agent -> profile",  MID_BLUE),
    ("5",  "Parallel Research",       "Weather + Transport + Hotel + Places all at once  (saves ~45s)",RGBColor(0x00, 0x7A, 0xCC)),
    ("6",  "Conflict Detection",      "Orchestrator detects issues, retries specific agents (max 3x)", ORANGE),
    ("7",  "Budget + Itinerary",      "Deterministic cost calc  ->  LLM day-by-day schedule",         ACCENT),
    ("8",  "Final Review",            "LLM QA gate -> approve / retry once / approve anyway",          ORANGE),
    ("9",  "Output Guardrails",       "6 post-generation validators before PDF creation",              RED),
    ("10", "PDF Generation",          "ReportLab 8-section report  +  ChromaDB memory update",        PURPLE),
    ("11", "Evaluation Suite",        "Basic evals (free)  +  RAGAS  +  DeepEval  ->  JSON report",   TEAL),
    ("12", "Streamlit shows results", "Progress bar hits 100%  ->  tabs + PDF download appear",        ACCENT),
]

# 4 cols x 3 rows
CW2 = (12.83 - 3 * 0.10) / 4   # = 3.13
CH2 = (CONTENT_B - CONTENT_Y - 2 * 0.08) / 3   # = 1.83

for i, (num, title, detail, bg_) in enumerate(steps):
    col = i % 4
    row = i // 4
    x = MARGIN + col * (CW2 + 0.10)
    y = CONTENT_Y + row * (CH2 + 0.08)

    rect(s, x, y, CW2, CH2, fill=WHITE, line=bg_, lw=Pt(1.5))
    rect(s, x, y, 0.52, CH2, fill=bg_)
    txt(s, num, x, y, 0.52, CH2, size=15, bold=True,
        color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, title,  x + 0.60, y + 0.12, CW2 - 0.68, 0.36, size=11, bold=True, color=DARK_BLUE)
    txt(s, detail, x + 0.60, y + 0.52, CW2 - 0.68, CH2 - 0.60, size=10, color=DARK_GREY)

    if col < 3:
        txt(s, ">", x + CW2 + 0.01, y + 0.62, 0.09, 0.50,
            size=12, bold=True, color=MID_BLUE, align=PP_ALIGN.CENTER)

footer(s, "Total: ~60-120s  |  Parallel research saves ~45s  |  Guardrails + evals add ~5s  |  Everything stored before user sees result")


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 18 — KEY DESIGN DECISIONS
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header(s, "Key Design Decisions", "What makes this system stand out")

features = [
    ("Parallel Execution",       MID_BLUE,
     "LangGraph Send() fans out 4 agents simultaneously.\nSaves ~45s. Auto-synced via merge_dicts reducer."),
    ("Structured LLM Output",    DARK_BLUE,
     ".with_structured_output(Pydantic) forces typed JSON.\nNo manual parsing. Bad outputs rejected automatically."),
    ("3-Layer Input Guardrails", RED,
     "Rules -> LLM safety -> PII masking — in order.\nBlocks before any LLM cost is ever incurred."),
    ("6-Check Output Guardrails",RED,
     "Per-agent post-generation validators.\nCatches bad data before downstream agents use it."),
    ("No LLM for Math",          ACCENT,
     "Budget Agent is pure Python. Faster, cheaper, exact.\nLLMs hallucinate numbers — Python math does not."),
    ("Semantic Memory",          TEAL,
     "ChromaDB stores past trips as vectors.\nReturning users get personalised recommendations."),
    ("Conflict-Driven Retries",  ORANGE,
     "Orchestrator detects issues and retries agents.\nNo fixed pipeline — adapts to what went wrong."),
    ("Evaluation Suite",         PURPLE,
     "Basic (free, instant) + RAGAS + DeepEval.\nSurfaced real bugs. 76 tests, all passing."),
    ("Graceful Degradation",     DARK_GREY,
     "No API key? Mock data. Search fails? Fallbacks.\nUser always gets a plan — never a blank error."),
]

# 3 cols x 3 rows
CW3 = (12.83 - 2 * 0.18) / 3   # = 4.16
CH3 = (CONTENT_B - CONTENT_Y - 2 * 0.10) / 3   # = 1.81

for i, (title, bg_, body) in enumerate(features):
    col = i % 3
    row = i // 3
    x = MARGIN + col * (CW3 + 0.18)
    y = CONTENT_Y + row * (CH3 + 0.10)
    rect(s, x, y, CW3, CH3, fill=WHITE, line=bg_, lw=Pt(2))
    rect(s, x, y, CW3, 0.44, fill=bg_)
    txt(s, title, x + 0.14, y + 0.07, CW3 - 0.22, 0.32, size=13, bold=True, color=WHITE)
    txt(s, body,  x + 0.14, y + 0.54, CW3 - 0.22, CH3 - 0.62, size=11, color=DARK_GREY)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 19 — LIVE DEMO
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, 13.33, 7.50, fill=DARK_BLUE)
rect(s, 0, 0, 13.33, 0.07, fill=MID_BLUE)
rect(s, 0, 7.43, 13.33, 0.07, fill=MID_BLUE)

# Decorative circles — kept inside slide bounds
for cx, cy, r, c in [
    (10.0, 0.5, 2.2, MID_BLUE),
    (11.5, 2.2, 1.3, RGBColor(0x0D, 0x2A, 0x45)),
]:
    sh = s.shapes.add_shape(9, Inches(cx), Inches(cy), Inches(r), Inches(r))
    sh.fill.solid(); sh.fill.fore_color.rgb = c; sh.line.fill.background()

txt(s, "LIVE DEMO", 0, 0.9, 13.33, 1.3,
    size=70, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
rect(s, 2.5, 2.35, 8.33, 0.06, fill=ACCENT)

for i, (action, detail, url) in enumerate([
    ("Streamlit UI",        "Open in your browser",                    "http://localhost:8501"),
    ("Fill trip details",   "Origin, destination, dates, budget",      ""),
    ("Click Plan My Trip",  "Watch phase-by-phase progress bar live",  ""),
    ("View results",        "Itinerary / Budget / Transport / Review", ""),
    ("Download PDF",        "8-section professional travel report",    ""),
    ("API Docs",            "All endpoints + request/response models", "http://localhost:8000/docs"),
]):
    y = 2.52 + i * 0.74
    rect(s, 2.50, y, 0.54, 0.58, fill=MID_BLUE)
    txt(s, str(i + 1), 2.50, y, 0.54, 0.58, size=15, bold=True,
        color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, action, 3.14, y + 0.05, 2.80, 0.28, size=13, bold=True, color=YELLOW)
    txt(s, detail, 3.14, y + 0.34, 3.30, 0.22, size=11,
        color=RGBColor(0xB8, 0xD4, 0xF0))
    if url:
        txt(s, url, 6.55, y + 0.12, 4.00, 0.28, size=12, bold=True, color=ACCENT)

# pytest result strip — inside slide bounds
rect(s, 1.50, 7.00, 10.33, 0.36, fill=MID_BLUE)
txt(s, "pytest tests/  ->  76 passed in 2.3 seconds", 1.50, 7.02, 10.33, 0.28,
    size=13, bold=True, color=YELLOW, align=PP_ALIGN.CENTER)


# ── Save ──────────────────────────────────────────────────────────────────────
OUT = r"c:\Users\praveen.kg\Desktop\Praveen\Practice\trip_planner\AI_Trip_Planner.pptx"
prs.save(OUT)
print(f"Saved  ->  {OUT}")
print(f"Slides :   {len(prs.slides)}")
